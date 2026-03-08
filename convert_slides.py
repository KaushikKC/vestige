#!/usr/bin/env python3
"""
Convert vestige-pitch.html → vestige-pitch.pdf + vestige-pitch.pptx
Uses:
  - Chrome headless for PDF (16:9, no margins)
  - python-pptx for PPTX (pixel-perfect recreation with Vestige theme)
"""

import subprocess, os, sys, time
from pathlib import Path

# ── Paths ──────────────────────────────────────────────────────────────
BASE      = Path(__file__).parent
HTML_SRC  = BASE / "vestige-pitch.html"
PRINT_HTML= BASE / "vestige-print.html"
PDF_OUT   = BASE / "vestige-pitch.pdf"
PPTX_OUT  = BASE / "vestige-pitch.pptx"
CHROME    = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"


# ══════════════════════════════════════════════════════════════════════
# STEP 1 — Build a print-optimised HTML (16:9 pages, reveals visible)
# ══════════════════════════════════════════════════════════════════════
def build_print_html():
    src = HTML_SRC.read_text(encoding="utf-8")

    print_css = """
<style id="print-override">
@page {
  size: 13.33in 7.5in;
  margin: 0;
}
@media print {
  html { scroll-snap-type: none !important; scroll-behavior: auto !important; }
  body { overflow: visible !important; }
  .slide {
    width: 100vw !important;
    height: 100vh !important;
    min-height: 0 !important;
    overflow: hidden !important;
    page-break-after: always !important;
    break-after: page !important;
    position: relative !important;
  }
  .slide:last-child { page-break-after: avoid !important; break-after: avoid !important; }
  .slide-content { overflow: hidden !important; max-height: 100% !important; }
  /* Show all animated elements immediately */
  .reveal, .reveal-left, .reveal-scale {
    opacity: 1 !important;
    transform: none !important;
    transition: none !important;
  }
  /* Hide interactive chrome */
  .progress-bar, .nav-dots, .slide-counter,
  .keyboard-hint, .edit-hotzone, .edit-toggle,
  .edit-banner { display: none !important; }
  /* Stop CSS animations */
  * { animation: none !important; }
}
</style>
"""
    # Insert just before </head>
    out = src.replace("</head>", print_css + "\n</head>")
    PRINT_HTML.write_text(out, encoding="utf-8")
    print(f"  ✓ Print HTML written: {PRINT_HTML.name}")


# ══════════════════════════════════════════════════════════════════════
# STEP 2 — Chrome headless → PDF
# ══════════════════════════════════════════════════════════════════════
def build_pdf():
    url = PRINT_HTML.as_uri()
    cmd = [
        CHROME,
        "--headless=new",
        "--disable-gpu",
        "--no-sandbox",
        "--disable-dev-shm-usage",
        "--run-all-compositor-stages-before-draw",
        "--virtual-time-budget=4000",
        f"--print-to-pdf={PDF_OUT}",
        "--print-to-pdf-no-header",
        "--no-pdf-header-footer",
        url,
    ]
    print(f"  Running Chrome headless…")
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
    if PDF_OUT.exists():
        size_kb = PDF_OUT.stat().st_size // 1024
        print(f"  ✓ PDF written: {PDF_OUT.name}  ({size_kb} KB)")
    else:
        print("  ✗ PDF generation failed")
        print(result.stderr[-800:] if result.stderr else "(no stderr)")
        sys.exit(1)


# ══════════════════════════════════════════════════════════════════════
# STEP 3 — python-pptx → PPTX
# Recreates all 9 slides with Vestige theme colours & typography.
# ══════════════════════════════════════════════════════════════════════
def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt
    import pptx.oxml.ns as nsmap
    from lxml import etree

    # ── Vestige palette (from theme.ts) ─────────────────────────────
    C_BG       = RGBColor(0x0C, 0x0D, 0x10)
    C_SURFACE  = RGBColor(0x11, 0x12, 0x16)
    C_CARD     = RGBColor(0x17, 0x18, 0x1D)
    C_ACCENT   = RGBColor(0xF5, 0xF1, 0x00)
    C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
    C_SECONDARY= RGBColor(0xB3, 0xB3, 0xB8)
    C_TERTIARY = RGBColor(0x6E, 0x6E, 0x73)
    C_DIVIDER  = RGBColor(0x23, 0x24, 0x2A)
    C_SUCCESS  = RGBColor(0x22, 0xC5, 0x5E)
    C_ERROR    = RGBColor(0xEF, 0x44, 0x44)

    # ── Slide dimensions: 16:9 widescreen ───────────────────────────
    W = Inches(13.33)
    H = Inches(7.5)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # Helper: blank slide layout
    blank = prs.slide_layouts[6]  # completely blank

    # ── Low-level helpers ────────────────────────────────────────────
    def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_width_pt=0):
        shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE.RECTANGLE
        fill = shape.fill
        if fill_rgb:
            fill.solid()
            fill.fore_color.rgb = fill_rgb
        else:
            fill.background()
        line = shape.line
        if line_rgb:
            line.color.rgb = line_rgb
            line.width = Pt(line_width_pt)
        else:
            line.fill.background()
        return shape

    def add_text(slide, text, x, y, w, h,
                 font_size=18, bold=False, color=None,
                 align=PP_ALIGN.LEFT, wrap=True,
                 italic=False, letter_spacing=0):
        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color or C_WHITE
        run.font.name = "Space Grotesk"
        return txBox

    def add_accent_bar(slide, x, y, w=Inches(0.55), h=Pt(3.5)):
        bar = add_rect(slide, x, y, w, h, fill_rgb=C_ACCENT)
        bar.line.fill.background()
        return bar

    def slide_bg(slide):
        """Fill slide background with C_BG and add subtle grid lines."""
        add_rect(slide, 0, 0, W, H, fill_rgb=C_BG)
        # Subtle grid — 4 horizontal + 4 vertical lines
        grid_color = RGBColor(0x1A, 0x1B, 0x20)
        for i in range(1, 5):
            # Horizontal
            ln = add_rect(slide, 0, H * i // 5, W, Pt(0.5), fill_rgb=grid_color)
            # Vertical
            ln2= add_rect(slide, W * i // 5, 0, Pt(0.5), H, fill_rgb=grid_color)

    def section_label(slide, text, x, y):
        add_text(slide, text.upper(), x, y, Inches(5), Inches(0.3),
                 font_size=9, bold=True, color=C_TERTIARY, letter_spacing=2)

    def card(slide, x, y, w, h):
        r = add_rect(slide, x, y, w, h, fill_rgb=C_CARD, line_rgb=C_DIVIDER, line_width_pt=0.75)
        return r

    # ════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title
    # ════════════════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(blank)
    slide_bg(s1)

    # Hackathon chip
    chip = add_rect(s1, Inches(1), Inches(1.1), Inches(3.2), Inches(0.32),
                    fill_rgb=RGBColor(0x1A, 0x1A, 0x0A), line_rgb=RGBColor(0x50, 0x4A, 0x00), line_width_pt=0.75)
    add_text(s1, "SOLANA MOBILE HACKATHON 2025",
             Inches(1.05), Inches(1.12), Inches(3.1), Inches(0.28),
             font_size=7.5, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # VESTIGE wordmark — "V" in yellow, rest white
    add_text(s1, "V", Inches(1), Inches(1.55), Inches(0.85), Inches(1.5),
             font_size=96, bold=True, color=C_ACCENT)
    add_text(s1, "ESTIGE", Inches(1.72), Inches(1.55), Inches(4.5), Inches(1.5),
             font_size=96, bold=True, color=C_WHITE)

    # Accent bar
    add_accent_bar(s1, Inches(1), Inches(3.2))

    # Tagline
    add_text(s1, "Fair-launch token trading with an inverted bonding curve\n— built natively for Solana Mobile.",
             Inches(1), Inches(3.45), Inches(5.5), Inches(1),
             font_size=15, color=C_SECONDARY, wrap=True)

    # Chips row
    chips = ["Solana", "DeFi", "Mobile-First", "Raydium CPMM"]
    cx = Inches(1)
    for ch in chips:
        w_chip = Inches(max(0.7, len(ch) * 0.095 + 0.25))
        add_rect(s1, cx, Inches(4.65), w_chip, Inches(0.28),
                 fill_rgb=C_SURFACE, line_rgb=C_DIVIDER, line_width_pt=0.75)
        add_text(s1, ch, cx + Inches(0.06), Inches(4.67), w_chip - Inches(0.1), Inches(0.25),
                 font_size=8.5, color=C_SECONDARY, align=PP_ALIGN.CENTER)
        cx += w_chip + Inches(0.12)

    # Decorative rings (right side) — concentric circles
    cx_ring, cy_ring, r_ring = Inches(10.2), Inches(3.75), Inches(2.2)
    for i, (alpha, size) in enumerate([(0.12, 1.0), (0.08, 0.76), (0.18, 0.52)]):
        rc = int(0xF5 * alpha)
        gc = int(0xF1 * alpha)
        bc = int(0x00 * alpha)
        rad = r_ring * size
        circ = s1.shapes.add_shape(9, cx_ring - rad, cy_ring - rad, rad*2, rad*2)  # oval
        circ.fill.background()
        circ.line.color.rgb = RGBColor(0x35, 0x33, 0x00) if alpha > 0.1 else RGBColor(0x25, 0x24, 0x00)
        circ.line.width = Pt(1)


    # ════════════════════════════════════════════════════════════════
    # SLIDE 2 — The Problem
    # ════════════════════════════════════════════════════════════════
    s2 = prs.slides.add_slide(blank)
    slide_bg(s2)

    add_accent_bar(s2, Inches(1), Inches(0.9))
    section_label(s2, "The Problem", Inches(1), Inches(1.08))
    add_text(s2, "Token launches are broken.", Inches(1), Inches(1.45), Inches(8), Inches(0.9),
             font_size=44, bold=True, color=C_WHITE)
    add_text(s2, "Insider-heavy launches, bot-driven rug pulls, and zero mobile UX\nleave retail traders with no fair entry point.",
             Inches(1), Inches(2.45), Inches(7), Inches(0.8),
             font_size=14, color=C_SECONDARY, wrap=True)

    stat_data = [
        ("~80%", "of meme coins abandoned\nwithin 48 hours", C_ERROR),
        ("$2B+",  "lost to rug pulls &\ninsider dumps in 2024", RGBColor(0xFF, 0x9F, 0x0A)),
        ("0",     "native Solana Mobile\ntoken-launch apps exist", C_TERTIARY),
    ]
    for i, (num, label, color) in enumerate(stat_data):
        sx = Inches(1 + i * 3.8)
        card(s2, sx, Inches(3.45), Inches(3.5), Inches(2.9))
        add_text(s2, num, sx + Inches(0.22), Inches(3.65), Inches(3), Inches(1),
                 font_size=52, bold=True, color=color)
        add_text(s2, label, sx + Inches(0.22), Inches(4.8), Inches(3), Inches(0.85),
                 font_size=11.5, color=C_SECONDARY, wrap=True)


    # ════════════════════════════════════════════════════════════════
    # SLIDE 3 — The Solution
    # ════════════════════════════════════════════════════════════════
    s3 = prs.slides.add_slide(blank)
    slide_bg(s3)

    add_accent_bar(s3, Inches(1), Inches(0.9))
    section_label(s3, "The Solution", Inches(1), Inches(1.08))
    add_text(s3, "An inverted bonding curve\nthat rewards early holders.",
             Inches(1), Inches(1.45), Inches(6.2), Inches(1.6),
             font_size=36, bold=True, color=C_WHITE, wrap=True)

    add_text(s3, "Vestige starts tokens at their maximum price and decreases it as supply "
                 "is sold — creating a built-in incentive for early participation and "
                 "eliminating the insider advantage.",
             Inches(1), Inches(3.2), Inches(5.8), Inches(1),
             font_size=12.5, color=C_SECONDARY, wrap=True)

    add_text(s3, "Once the SOL graduation target is hit, the token graduates automatically "
                 "to Raydium CPMM — giving it real on-chain liquidity.",
             Inches(1), Inches(4.35), Inches(5.8), Inches(0.85),
             font_size=12.5, color=C_SECONDARY, wrap=True)

    # Curve diagram box
    card(s3, Inches(7.6), Inches(1.0), Inches(5.0), Inches(5.6))
    add_text(s3, "Price", Inches(7.75), Inches(1.1), Inches(1), Inches(0.3),
             font_size=9, color=C_TERTIARY)
    add_text(s3, "↑  pMax", Inches(7.75), Inches(1.35), Inches(1.5), Inches(0.3),
             font_size=8.5, color=C_SECONDARY)
    add_text(s3, "pMin", Inches(11.2), Inches(5.8), Inches(1), Inches(0.3),
             font_size=8.5, color=C_SECONDARY)
    add_text(s3, "→ Token Supply Sold", Inches(9.2), Inches(6.1), Inches(3), Inches(0.3),
             font_size=9, color=C_TERTIARY)
    add_text(s3, "→ Raydium CPMM at graduation", Inches(10.0), Inches(4.5), Inches(2.5), Inches(0.5),
             font_size=8, color=C_SUCCESS, wrap=True)

    # Draw inverted curve as a connector
    from pptx.util import Emu
    conn = s3.shapes.add_connector(1,
        Inches(8.1), Inches(1.55),   # start (top-left = high price)
        Inches(12.1), Inches(5.75))  # end (bottom-right = low price)
    conn.line.color.rgb = C_ACCENT
    conn.line.width = Pt(2.5)


    # ════════════════════════════════════════════════════════════════
    # SLIDE 4 — How It Works
    # ════════════════════════════════════════════════════════════════
    s4 = prs.slides.add_slide(blank)
    slide_bg(s4)

    add_accent_bar(s4, Inches(1), Inches(0.9))
    section_label(s4, "Mechanism", Inches(1), Inches(1.08))
    add_text(s4, "How It Works", Inches(1), Inches(1.45), Inches(8), Inches(0.75),
             font_size=42, bold=True, color=C_WHITE)

    steps = [
        ("01", "Creator Launches a Token",
         "Sets name, symbol, supply, and graduation SOL target. Token starts at pMax.", "Launch", C_ACCENT),
        ("02", "Traders Buy & Sell via Bonding Curve",
         "price = pMax − (pMax − pMin) × (totalBaseSold / supply). Every buy lowers cost for next buyer.", "Buy / Sell", C_SUCCESS),
        ("03", "Live Price Feed & Trade History",
         "Candlestick chart, bonding curve visualization, and real-time trade feed parsed from on-chain logs.", "Charts", C_SECONDARY),
        ("04", "Automatic Raydium Graduation",
         "When totalSolCollected ≥ graduationTarget, liquidity is seeded into Raydium CPMM automatically.", "Graduate", C_ACCENT),
    ]
    for i, (num, title, desc, badge, badge_color) in enumerate(steps):
        sy = Inches(2.45 + i * 1.1)
        card(s4, Inches(1), sy, Inches(11.3), Inches(0.95))
        add_text(s4, num, Inches(1.18), sy + Inches(0.18), Inches(0.5), Inches(0.65),
                 font_size=13, bold=True, color=C_ACCENT)
        add_text(s4, title, Inches(1.75), sy + Inches(0.08), Inches(6.5), Inches(0.35),
                 font_size=13, bold=True, color=C_WHITE)
        add_text(s4, desc, Inches(1.75), sy + Inches(0.47), Inches(7.5), Inches(0.42),
                 font_size=10, color=C_SECONDARY, wrap=True)
        # Badge
        bw = Inches(max(0.7, len(badge) * 0.1 + 0.2))
        add_rect(s4, Inches(11.5) - bw, sy + Inches(0.28), bw, Inches(0.32),
                 fill_rgb=RGBColor(0x1A, 0x1A, 0x0E) if badge_color == C_ACCENT else C_CARD,
                 line_rgb=C_DIVIDER, line_width_pt=0.5)
        add_text(s4, badge, Inches(11.5) - bw + Inches(0.05), sy + Inches(0.29),
                 bw - Inches(0.08), Inches(0.3),
                 font_size=8, bold=True, color=badge_color, align=PP_ALIGN.CENTER)


    # ════════════════════════════════════════════════════════════════
    # SLIDE 5 — Product Features
    # ════════════════════════════════════════════════════════════════
    s5 = prs.slides.add_slide(blank)
    slide_bg(s5)

    add_accent_bar(s5, Inches(1), Inches(0.75))
    section_label(s5, "Product", Inches(1), Inches(0.93))
    add_text(s5, "Built for traders on the go", Inches(1), Inches(1.25), Inches(8), Inches(0.65),
             font_size=38, bold=True, color=C_WHITE)

    feats = [
        ("📈", "Bonding Curve Chart",   "Real-time SVG of the inverted curve with a live position dot."),
        ("🕯️", "Candlestick Charts",    "Custom OHLC candles aggregated from on-chain transaction logs."),
        ("⚡", "Live Trade Feed",       "Buy/sell events streamed from Solana logs in real time."),
        ("🚀", "One-Tap Launch",        "Create and launch a token in seconds from a Solana Mobile device."),
        ("🎓", "Auto Graduation",       "Automatic Raydium CPMM seeding when the SOL target is reached."),
        ("📱", "Mobile-Native",         "React Native + Solana Mobile SDK — hardware wallet signing on-device."),
    ]
    cols, rows = 3, 2
    for i, (icon, name, desc) in enumerate(feats):
        col = i % cols
        row = i // cols
        fx = Inches(1 + col * 3.9)
        fy = Inches(2.1 + row * 2.35)
        card(s5, fx, fy, Inches(3.65), Inches(2.1))
        # Icon box
        add_rect(s5, fx + Inches(0.18), fy + Inches(0.18), Inches(0.45), Inches(0.45),
                 fill_rgb=RGBColor(0x22, 0x20, 0x00), line_rgb=RGBColor(0x40, 0x3C, 0x00), line_width_pt=0.5)
        add_text(s5, icon, fx + Inches(0.19), fy + Inches(0.17), Inches(0.44), Inches(0.44),
                 font_size=14, align=PP_ALIGN.CENTER)
        add_text(s5, name, fx + Inches(0.18), fy + Inches(0.72), Inches(3.3), Inches(0.35),
                 font_size=12.5, bold=True, color=C_WHITE)
        add_text(s5, desc, fx + Inches(0.18), fy + Inches(1.1), Inches(3.3), Inches(0.85),
                 font_size=10.5, color=C_SECONDARY, wrap=True)


    # ════════════════════════════════════════════════════════════════
    # SLIDE 6 — Tech Stack
    # ════════════════════════════════════════════════════════════════
    s6 = prs.slides.add_slide(blank)
    slide_bg(s6)

    add_accent_bar(s6, Inches(1), Inches(0.75))
    section_label(s6, "Architecture", Inches(1), Inches(0.93))
    add_text(s6, "Fully on-chain, fully mobile", Inches(1), Inches(1.25), Inches(8), Inches(0.65),
             font_size=38, bold=True, color=C_WHITE)

    # Two columns
    for col_i, (col_title, items, dot_color) in enumerate([
        ("On-Chain (Solana)", [
            "Anchor framework — Rust programs",
            "Custom inverted bonding curve logic",
            "Raydium CPMM graduation CPI",
            "On-chain log emission for trade events",
            "BN arithmetic for lamport precision",
        ], C_ACCENT),
        ("Mobile (React Native)", [
            "Solana Mobile SDK (MWA)",
            "React Native + TypeScript",
            "Custom SVG candle + curve charts",
            "Real-time log parsing & candle aggregation",
            "Auto-refresh every 15 / 30 seconds",
        ], C_SUCCESS),
    ]):
        cx = Inches(1 + col_i * 6.2)
        add_rect(s6, cx, Inches(2.15), Inches(5.8), Pt(1), fill_rgb=C_DIVIDER)
        add_text(s6, col_title, cx, Inches(2.05), Inches(5.5), Inches(0.38),
                 font_size=14, bold=True, color=C_WHITE)
        for j, item in enumerate(items):
            iy = Inches(2.65 + j * 0.82)
            # dot
            add_rect(s6, cx + Inches(0.05), iy + Inches(0.14), Inches(0.09), Inches(0.09),
                     fill_rgb=dot_color)
            add_text(s6, item, cx + Inches(0.22), iy, Inches(5.4), Inches(0.55),
                     font_size=12, color=C_SECONDARY, wrap=True)


    # ════════════════════════════════════════════════════════════════
    # SLIDE 7 — Market Opportunity
    # ════════════════════════════════════════════════════════════════
    s7 = prs.slides.add_slide(blank)
    slide_bg(s7)

    add_accent_bar(s7, Inches(1), Inches(0.75))
    section_label(s7, "Opportunity", Inches(1), Inches(0.93))
    add_text(s7, "A massive, underserved market", Inches(1), Inches(1.25), Inches(8), Inches(0.65),
             font_size=38, bold=True, color=C_WHITE)

    markets = [
        ("SOLANA DEX VOLUME", "$40B+",  "Monthly on-chain volume\nacross Solana DEXs in 2025"),
        ("PUMP.FUN LAUNCHES",  "4M+",   "Tokens launched since 2024\n— demand is proven"),
        ("SOLANA MOBILE USERS","~200K", "Saga + Chapter 2 owners\n— zero native launch apps"),
    ]
    for i, (title, fig, sub) in enumerate(markets):
        mx = Inches(1 + i * 4.1)
        card(s7, mx, Inches(2.25), Inches(3.8), Inches(2.9))
        add_text(s7, title, mx + Inches(0.22), Inches(2.42), Inches(3.4), Inches(0.32),
                 font_size=8, bold=True, color=C_TERTIARY, letter_spacing=2)
        add_text(s7, fig, mx + Inches(0.22), Inches(2.78), Inches(3.4), Inches(1),
                 font_size=52, bold=True, color=C_ACCENT)
        add_text(s7, sub, mx + Inches(0.22), Inches(3.95), Inches(3.3), Inches(0.85),
                 font_size=11, color=C_SECONDARY, wrap=True)

    add_text(s7,
        "Vestige captures the intersection of fair-launch demand, Solana DeFi growth, "
        "and the untapped Solana Mobile ecosystem — a vertical no competitor targets natively.",
        Inches(1), Inches(5.4), Inches(11.3), Inches(0.75),
        font_size=12.5, color=C_SECONDARY, wrap=True)


    # ════════════════════════════════════════════════════════════════
    # SLIDE 8 — Roadmap
    # ════════════════════════════════════════════════════════════════
    s8 = prs.slides.add_slide(blank)
    slide_bg(s8)

    add_accent_bar(s8, Inches(1), Inches(0.75))
    section_label(s8, "Roadmap", Inches(1), Inches(0.93))
    add_text(s8, "What's next for Vestige", Inches(1), Inches(1.25), Inches(8), Inches(0.65),
             font_size=38, bold=True, color=C_WHITE)

    phases = [
        ("✓ Complete", "Foundation", C_SUCCESS,
         ["Anchor bonding curve program", "Mobile buy / sell UI",
          "Price + candle charts", "Raydium graduation CPI"]),
        ("▶ Now", "Hackathon MVP", C_ACCENT,
         ["Live trade feed", "Create launch screen",
          "Portfolio view", "Mainnet beta"]),
        ("Q3 2025", "Growth", C_TERTIARY,
         ["Creator analytics", "Social trading feed",
          "Push notifications", "Token metadata CDN"]),
        ("Q4 2025", "Scale", C_TERTIARY,
         ["Multi-chain support", "Revenue share for creators",
          "SDK for third-party launchpads", "DAO governance"]),
    ]
    for i, (label, name, label_color, items) in enumerate(phases):
        px = Inches(1 + i * 3.1)
        is_active = label == "▶ Now"
        bg = RGBColor(0x1C, 0x1B, 0x08) if is_active else C_CARD
        bord = RGBColor(0x50, 0x4C, 0x10) if is_active else C_DIVIDER
        card_h = Inches(4.6)
        add_rect(s8, px, Inches(2.2), Inches(2.85), card_h, fill_rgb=bg, line_rgb=bord, line_width_pt=0.75)
        add_text(s8, label, px + Inches(0.18), Inches(2.35), Inches(2.5), Inches(0.32),
                 font_size=9, bold=True, color=label_color)
        add_text(s8, name, px + Inches(0.18), Inches(2.68), Inches(2.5), Inches(0.38),
                 font_size=14, bold=True, color=C_WHITE)
        for j, item in enumerate(items):
            add_text(s8, "→  " + item, px + Inches(0.18), Inches(3.18 + j * 0.72), Inches(2.6), Inches(0.62),
                     font_size=10, color=C_SECONDARY, wrap=True)


    # ════════════════════════════════════════════════════════════════
    # SLIDE 9 — The Ask
    # ════════════════════════════════════════════════════════════════
    s9 = prs.slides.add_slide(blank)
    slide_bg(s9)
    # Glow background circle
    glow = s9.shapes.add_shape(9, Inches(3.66), Inches(0.75), Inches(6), Inches(6))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(0x12, 0x11, 0x00)
    glow.line.fill.background()

    add_accent_bar(s9, Inches(6.39), Inches(0.85), w=Inches(0.55))

    section_label(s9, "Seed Round", Inches(5.5), Inches(1.1))

    add_text(s9, "$500K", Inches(3.0), Inches(1.4), Inches(7.3), Inches(1.8),
             font_size=108, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    add_text(s9,
        "to take Vestige from hackathon MVP to the #1 fair-launch app on Solana Mobile.",
        Inches(2.5), Inches(3.25), Inches(8.3), Inches(0.7),
        font_size=14, color=C_SECONDARY, align=PP_ALIGN.CENTER, wrap=True)

    use_data = [("40%", "Engineering\n2 senior Solana devs"),
                ("30%", "Growth &\ncommunity incentives"),
                ("30%", "Security audits\n& infrastructure")]
    for i, (pct, desc) in enumerate(use_data):
        ux = Inches(1.55 + i * 3.5)
        card(s9, ux, Inches(4.15), Inches(3.2), Inches(1.65))
        add_text(s9, pct, ux + Inches(0.18), Inches(4.28), Inches(3), Inches(0.45),
                 font_size=22, bold=True, color=C_WHITE)
        add_text(s9, desc, ux + Inches(0.18), Inches(4.75), Inches(2.85), Inches(0.85),
                 font_size=10.5, color=C_SECONDARY, wrap=True)

    contacts = ["vestige.app", "@vestige_sol", "team@vestige.app"]
    for i, c in enumerate(contacts):
        cx2 = Inches(2.6 + i * 2.8)
        add_rect(s9, cx2, Inches(6.0), Inches(2.4), Inches(0.3),
                 fill_rgb=C_SURFACE, line_rgb=C_DIVIDER, line_width_pt=0.5)
        add_text(s9, c, cx2 + Inches(0.08), Inches(6.01), Inches(2.25), Inches(0.28),
                 font_size=9.5, color=C_TERTIARY, align=PP_ALIGN.CENTER)

    # ── Save PPTX ───────────────────────────────────────────────────
    prs.save(PPTX_OUT)
    size_kb = PPTX_OUT.stat().st_size // 1024
    print(f"  ✓ PPTX written: {PPTX_OUT.name}  ({size_kb} KB)")


# ══════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    print("\n🎞  Vestige Slide Converter\n")

    print("Step 1/3  Building print HTML…")
    build_print_html()

    print("Step 2/3  Exporting PDF via Chrome headless…")
    build_pdf()

    print("Step 3/3  Building PPTX via python-pptx…")
    build_pptx()

    # Clean up temp file
    if PRINT_HTML.exists():
        PRINT_HTML.unlink()
        print(f"\n  ✓ Cleaned up temp file")

    print(f"\n✅  Done!\n   PDF  → {PDF_OUT}\n   PPTX → {PPTX_OUT}\n")
